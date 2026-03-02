import json
import pathlib
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Configuration
script_dir = pathlib.Path(__file__).parent.resolve()
# Input file is relative to the script location
INPUT_FILE = script_dir / "ai_capabilities.json"

# Theme: Midnight Executive
COLOR_PRIMARY = RGBColor(0x1E, 0x27, 0x61)    # Navy
COLOR_SECONDARY = RGBColor(0xCA, 0xDC, 0xFC)  # Ice Blue
COLOR_ACCENT = RGBColor(0xFF, 0xFF, 0xFF)     # White

def generate(output_dir):
    """
    Generates the AI Capabilities Matrix presentation.
    
    Args:
        output_dir: Directory to save the output file.
    """
    out_path = pathlib.Path(output_dir)
    OUTPUT_FILE = out_path / "AI_Capabilities_Matrix.pptx"

    print(f"Generating presentation from {INPUT_FILE} to {OUTPUT_FILE}...")

    # Load Data
    if not INPUT_FILE.exists():
        print(f"Error: {INPUT_FILE} not found.")
        return

    with open(INPUT_FILE, 'r') as f:
        data = json.load(f)

    prs = Presentation()
    
    # 1. Title Slide
    if 'title' in data:
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply Theme Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRIMARY

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = data.get('title', 'AI Capabilities Matrix')
        if title.text_frame:
            p = title.text_frame.paragraphs[0]
            p.font.color.rgb = COLOR_ACCENT
            p.font.bold = True
        
        subtitle.text = f"{data.get('subtitle', '')}\n{data.get('date', '')}"
        if subtitle.text_frame:
            p = subtitle.text_frame.paragraphs[0]
            p.font.color.rgb = COLOR_SECONDARY

    # 2. Matrix Slide
    if 'tools' in data and 'roles' in data:
        slide_layout = prs.slide_layouts[5] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Capability Matrix"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        # Table Data Preparation
        try:
            tools = data.get('tools', [])
            roles = data.get('roles', [])
            rows = len(roles) + 1
            cols = len(tools) + 1

            if rows > 1 and cols > 1:
                # Add Table
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9.0)
                height = Inches(4.0)
                
                table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                table = table_shape.table

                # Set Header (Tools)
                cell = table.cell(0, 0)
                cell.text = "Role / Tool"
                cell.text_frame.paragraphs[0].font.bold = True
                
                for i, tool in enumerate(tools):
                    cell = table.cell(0, i+1)
                    cell.text = tool
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_PRIMARY
                    if cell.text_frame.paragraphs:
                        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT

                # Set Body (Roles)
                for i, role in enumerate(roles):
                    # Row Header
                    cell = table.cell(i+1, 0)
                    cell.text = role.get('name', 'Unknown')
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_SECONDARY
                    if cell.text_frame.paragraphs:
                        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

                    # Values
                    evaluations = role.get('evaluations', {})
                    # If evaluations is a list (e.g. from JSON structure), convert to dict
                    if isinstance(evaluations, list):
                        evaluations = {e['tool']: e for e in evaluations}
                    
                    for j, tool in enumerate(tools):
                        # evaluations is now a dict: tool_name -> details
                        val_data = evaluations.get(tool, {"status": "N/A"})
                        # If simple string, wrap it
                        if isinstance(val_data, str):
                            val_data = {"status": val_data}
                            
                        status = val_data.get("status", "N/A")
                        note = val_data.get("note", "")
                        
                        cell = table.cell(i+1, j+1)
                        display_text = f"{status}\n({note})" if note else status
                        cell.text = display_text
                        cell.text_frame.paragraphs[0].font.size = Pt(9)
                    
        except Exception as e:
            print(f"Error accessing data keys or building table: {e}")

    # 3. Recommendations Slide
    if 'recommendations' in data:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Role Recommendations"
        if title.text_frame:
            title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

        content = slide.placeholders[1]
        tf = content.text_frame
        
        recommendations = data.get('recommendations', [])
        for rec in recommendations:
            if isinstance(rec, dict):
                p = tf.add_paragraph()
                role_val = rec.get('role', 'Role')
                tool_val = rec.get('tool', 'Tool')
                p.text = f"{role_val}: {tool_val}"
                p.font.bold = True
                p.font.size = Pt(24)
                p.font.color.rgb = COLOR_PRIMARY
                
                if 'reason' in rec:
                    p2 = tf.add_paragraph()
                    p2.text = f"Reason: {rec['reason']}"
                    p2.level = 1
                    p2.font.size = Pt(18)
            else:
                 # Fallback if just a string
                p = tf.add_paragraph()
                p.text = str(rec)


    # 4. Conclusion Slide (Dark Theme)
    if 'conclusion' in data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRIMARY
        
        title = slide.shapes.title
        conclusion_data = data['conclusion'] 
        points = []
        if isinstance(conclusion_data, dict):
            title.text = conclusion_data.get('title', 'Conclusion')
            points = conclusion_data.get('points', [])
        else:
            title.text = "Conclusion"
            # It might be a list of strings directly or a single string
            if isinstance(conclusion_data, list):
                points = conclusion_data
            else:
                points = [str(conclusion_data)]

        if title.text_frame:
            title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        for point in points:
            p = tf.add_paragraph()
            p.text = str(point)
            p.font.color.rgb = COLOR_SECONDARY
            p.font.size = Pt(24)

    # Ensure output directory exists
    out_path.mkdir(parents=True, exist_ok=True)
    
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    generate("outputs")
